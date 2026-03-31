"""
pptx_import.py — convert PPTX to slide images + extract notes
Usage: python pptx_import.py <pptx_path> <out_dir> [width] [height]
Output: JSON { images: [...], notes: [...] }

Windows: uses PowerPoint COM (win32com) — requires MS Office installed
Linux:   uses LibreOffice headless + pdf2image
"""
import sys, os, json, platform

def convert_windows(pptx_path, out_dir, width, height):
    import win32com.client
    pptx_path = os.path.abspath(pptx_path)
    out_dir   = os.path.abspath(out_dir)
    pptx_app  = win32com.client.Dispatch("PowerPoint.Application")
    pptx_app.Visible = True
    images = []
    try:
        prs = pptx_app.Presentations.Open(pptx_path, ReadOnly=True, Untitled=False, WithWindow=False)
        for i in range(1, prs.Slides.Count + 1):
            name = f"slide_{i:03d}.png"
            prs.Slides(i).Export(os.path.join(out_dir, name), "PNG", width, height)
            images.append(name)
        prs.Close()
    finally:
        pptx_app.Quit()
    return images

def convert_linux(pptx_path, out_dir, width, height):
    import subprocess, tempfile, glob
    from pdf2image import convert_from_path

    pptx_path = os.path.abspath(pptx_path)
    out_dir   = os.path.abspath(out_dir)

    with tempfile.TemporaryDirectory() as tmp:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            capture_output=True, text=True
        )
        pdfs = glob.glob(os.path.join(tmp, "*.pdf"))
        if not pdfs:
            raise RuntimeError(f"LibreOffice failed: {result.stderr.strip()}")

        pages = convert_from_path(pdfs[0], size=(width, height))
        images = []
        for i, page in enumerate(pages):
            name = f"slide_{i+1:03d}.png"
            page.save(os.path.join(out_dir, name), "PNG")
            images.append(name)
        return images

def extract_notes(pptx_path):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            notes.append(text)
        else:
            notes.append("")
    return notes

if __name__ == "__main__":
    pptx_path = sys.argv[1]
    out_dir   = sys.argv[2]
    width     = int(sys.argv[3]) if len(sys.argv) > 3 else 1920
    height    = int(sys.argv[4]) if len(sys.argv) > 4 else 1080

    if platform.system() == "Windows":
        images = convert_windows(pptx_path, out_dir, width, height)
    else:
        images = convert_linux(pptx_path, out_dir, width, height)

    notes = extract_notes(pptx_path)
    print(json.dumps({ "images": images, "notes": notes }))
